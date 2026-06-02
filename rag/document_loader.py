from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import json
import csv
import io


def load_document(uploaded_file):

    file_name = uploaded_file.name.lower()

    # PDF
    if file_name.endswith(".pdf"):

        reader = PdfReader(uploaded_file)

        text = ""

        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

        return text

    # DOCX
    elif file_name.endswith(".docx"):

        doc = Document(uploaded_file)

        return "\n".join(
            paragraph.text
            for paragraph in doc.paragraphs
        )

    # TXT / MD
    elif file_name.endswith((".txt", ".md")):

        return uploaded_file.read().decode(
            "utf-8",
            errors="ignore"
        )

    # JSON
    elif file_name.endswith(".json"):

        data = json.load(uploaded_file)

        return json.dumps(
            data,
            indent=2
        )

    # CSV
    elif file_name.endswith(".csv"):

        text = ""

        csv_content = uploaded_file.read().decode(
            "utf-8",
            errors="ignore"
        )

        reader = csv.reader(
            io.StringIO(csv_content)
        )

        for row in reader:
            text += " ".join(row) + "\n"

        return text

    # XLSX
    elif file_name.endswith(".xlsx"):

        workbook = load_workbook(
            uploaded_file,
            data_only=True
        )

        text = ""

        for sheet in workbook:

            for row in sheet.iter_rows(
                values_only=True
            ):

                text += " ".join(
                    str(cell)
                    for cell in row
                    if cell is not None
                )

                text += "\n"

        return text

    # PPTX
    elif file_name.endswith(".pptx"):

        presentation = Presentation(
            uploaded_file
        )

        text = ""

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text

    else:

        raise ValueError(
            f"Unsupported file type: {file_name}"
        )